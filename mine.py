#!/usr/bin/env python3
"""
FRCS (SN) Revision — Content Mining Orchestrator
=================================================
Tracks mining progress across all sources, extracts text chunks for Claude
to process into flashcards, validates content.js, and reports coverage gaps.

Usage:
  python3 mine.py status              # Show all sources and mining progress
  python3 mine.py next                # Show highest-priority unmined chunk
  python3 mine.py extract <id>        # Extract full text of a source
  python3 mine.py extract <id> 1-20  # Extract specific page range
  python3 mine.py done <id> 1-20 N   # Mark pages 1-20 as mined, N cards added
  python3 mine.py validate            # Check content.js for structural errors
  python3 mine.py stats               # Card counts by topic and korky tag
  python3 mine.py refs                # List all ref strings (spot missing pages)
"""

import json
import os
import re
import sys
from datetime import date
from pathlib import Path

# ── PATHS ────────────────────────────────────────────────────────────────────
APP_DIR   = Path.home() / "frcs-sn-revision"
KORKY_DIR = Path.home() / "Harvard University Dropbox/Carolyn Fu/Z. Personal/Sista & Me/FRCS Exam"
MANIFEST  = APP_DIR / "mining_manifest.json"
CONTENT   = APP_DIR / "content.js"

# ── SOURCE REGISTRY ──────────────────────────────────────────────────────────
# priority: 1 = mine first, 2 = high, 3 = medium, 4 = large reference
# type: recall_bank | landmark_paper | nice_guideline | course_material | textbook
SOURCES = {

    # ── RECALL BANK (priority 1 — real exam questions) ────────────────────
    "recall-questions": {
        "label": "Recall Bank — Questions (compiled)",
        "path": KORKY_DIR / "Recall Bank/Bank_Questions.pdf",
        "priority": 1, "type": "recall_bank",
        "topics": ["all"], "korky": True,
        "note": "Compiled recalled exam questions — highest fidelity for SBA generation",
    },
    "recall-2021": {
        "label": "Recall Bank 2021",
        "path": KORKY_DIR / "Recall Bank/Bank_Recall 2021.pdf",
        "priority": 1, "type": "recall_bank",
        "topics": ["all"], "korky": True,
        "note": "Recalled questions from 2021 sitting",
    },

    # ── LANDMARK PAPERS (priority 2) ──────────────────────────────────────
    "paper-isat": {
        "label": "ISAT — Molyneux 2009 (5-yr SAH outcomes)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2009 Molyneux ISAT 5-years Outcome.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["vascular-aneurysm"], "korky": True,
        "note": "ISAT: coiling vs clipping for ruptured intracranial aneurysms",
    },
    "paper-isuia": {
        "label": "ISUIA — Wiebers 1998 (unruptured aneurysms)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/1998 Wiebers ISUIA.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["vascular-aneurysm"], "korky": True,
        "note": "ISUIA: natural history and rupture risk of unruptured aneurysms",
    },
    "paper-stupp": {
        "label": "Stupp 2005 — RT + TMZ for GBM (NEJM)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2005 Stupp Radio plus Temozolomide for GBM.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["neuro-onco-cranial"], "korky": True,
        "note": "Stupp protocol: established 60Gy/30fr + concomitant/adjuvant TMZ as standard for GBM",
    },
    "paper-sport": {
        "label": "SPORT — Weinstein 2006 (spine RCT)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2006 Weinstein SPORT.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["degenerative-spine"], "korky": True,
        "note": "SPORT: surgery vs non-op for disc herniation and spinal stenosis",
    },
    "paper-decra": {
        "label": "DECRA — Cooper 2011 (decompressive craniectomy TBI)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2011 Cooper DECRA.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["head-injury"], "korky": True,
        "note": "DECRA: bifrontal DC for refractory ICP in diffuse TBI — worse outcomes",
    },
    "paper-destiny2": {
        "label": "DESTINY II — Juttler 2014 (malignant MCA, elderly)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2014 Juttler DESTINY II.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["vascular-ich"], "korky": True,
        "note": "DESTINY II: hemicraniectomy in patients >60 — survival benefit, disability caveat",
    },
    "paper-aruba": {
        "label": "ARUBA — Mohr 2014 (unruptured AVM)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2014 Mohr ARUBA.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["vascular-avm"], "korky": True,
        "note": "ARUBA: medical management superior to intervention for unruptured AVMs (controversial)",
    },
    "paper-nascis3": {
        "label": "NASCIS 3 — Bracken 1997 (steroids in SCI)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/1997 Bracken NASCIS 3.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["spinal-trauma"], "korky": True,
        "note": "NASCIS 3: methylprednisolone in acute SCI — now widely abandoned",
    },
    "paper-patchell": {
        "label": "Patchell 2005 — Surgery for malignant cord compression",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2005 Patchell Malig Cord Comp Surg.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["neuro-onco-spinal"], "korky": True,
        "note": "Patchell: surgery + RT superior to RT alone for MSCC — ability to walk primary endpoint",
    },
    "paper-berger": {
        "label": "Berger 2008 — Extent of resection LGG vs survival",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2008 Berger Extent of Resection of LGG vs Survival.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["neuro-onco-cranial"], "korky": True,
        "note": "Berger: EOR in LGG correlates with overall survival and malignant transformation",
    },
    "paper-santarius": {
        "label": "Santarius 2009 — CSDH drain (RCT)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2009 Santarius CSDH drain.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["head-injury"], "korky": True,
        "note": "Santarius: drain after burr-hole evacuation of CSDH reduces recurrence",
    },
    "paper-kulkarni": {
        "label": "Kulkarni 2010 — ETV Success Score",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2010 Kulkarni ETV Success Score.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["hydrocephalus", "paeds"], "korky": True,
        "note": "ETV Success Score: age + aetiology + prior shunt predict ETV success",
    },
    "paper-chesnut": {
        "label": "Chesnut 2012 — ICP monitoring in TBI (NEJM)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2012 Chesnut ICP Monitoring in TBI.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["head-injury", "neuro-icu"], "korky": True,
        "note": "Chesnut: ICP monitoring vs imaging-clinical exam — no mortality benefit in low/middle income settings",
    },
    "paper-liu-dbs": {
        "label": "Liu 2014 — GPi vs STN DBS for Parkinson's",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2014 Liu Metaanalysis of GPi vs STN for PD.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["functional"], "korky": True,
        "note": "GPi vs STN stimulation — comparable motor outcomes, different side-effect profiles",
    },
    "paper-hu-epilepsy": {
        "label": "Hu 2013 — SelAH vs ATL meta-analysis",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2013 Hu SelAH vs ATL Metaanalysis.pdf",
        "priority": 2, "type": "landmark_paper",
        "topics": ["epilepsy-surgery"], "korky": True,
        "note": "Selective amygdalohippocampectomy vs anterior temporal lobectomy — seizure outcomes",
    },
    "paper-mautner-nf2": {
        "label": "Mautner 2012 — Bevacizumab in NF2",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2012 Mautner Bevacizumab in NF2.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["paeds", "peripheral-nerve"], "korky": True,
        "note": "Bevacizumab for NF2-related vestibular schwannoma — hearing improvement",
    },
    "paper-casey-myelopathy": {
        "label": "Casey 1996 — Myelopathy surgery",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/1996 Casey RA Myelopathy Surgery.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["degenerative-spine"], "korky": True,
    },
    "paper-anand-mvd": {
        "label": "Anand 2011 — Safety of MVD in elderly",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2011 Anand Safe of MVD in Elderly.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["functional"], "korky": True,
        "note": "MVD for trigeminal neuralgia in elderly patients",
    },
    "paper-sivakumar-hyponatraemia": {
        "label": "Sivakumar 1994 — Hyponatraemia in Neurosurgery",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/1994 Sivakumar Management of Hyponatremia in Neurosurgery.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["neuro-icu"], "korky": True,
        "note": "SIADH vs CSW — frequently tested in recall bank",
    },
    "paper-smith-abscess": {
        "label": "Smith 2009 — Brain Abscess Surgery Review",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2009 Smith Abscess Surgery Review.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["neuro-onco-cranial", "neuro-icu"], "korky": True,
    },
    "paper-hashimoto-desh": {
        "label": "Hashimoto 2010 — DESH (NPH imaging)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2010 Hashimoto DESH.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["hydrocephalus"], "korky": True,
        "note": "Disproportionately enlarged subarachnoid spaces — NPH imaging criterion",
    },
    "paper-kumpe-iih-stents": {
        "label": "Kumpe 2011 — Dural Venous Stents for IIH",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2011 Kumpe Dural Stents for IIH.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["hydrocephalus"], "korky": True,
        "note": "Venous sinus stenting for IIH — multiple IIH questions in recall bank",
    },
    "paper-couture-helmet": {
        "label": "Couture 2013 — Helmet Therapy in 1050 Patients",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2013 Couture Helmet Therapy in 1050 Patients.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["paeds"], "korky": True,
        "note": "Plagiocephaly management — helmet vs repositioning",
    },
    "paper-delwel-nph": {
        "label": "Delwel 2013 — NPH Shunt Setting Trial (DEPSS)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2013 Delwel NPH Shunt Setting Trial DEPSS.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["hydrocephalus"], "korky": True,
        "note": "Dutch Evaluation of Programmable Shunt Settings — NPH shunt management",
    },
    "paper-jahangiri-redo-tss": {
        "label": "Jahangiri 2014 — Redo Transsphenoidal Complications",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Key Papers for FRCS/2014 Jahangiri Redo Transsphenoidal Complications.pdf",
        "priority": 3, "type": "landmark_paper",
        "topics": ["pituitary"], "korky": True,
        "note": "Complication rates in revision transsphenoidal surgery",
    },

    # ── RECALL BANK — additional files ────────────────────────────────────
    "recall-imp-topics": {
        "label": "Recall Bank — Important Topics & Questions",
        "path": KORKY_DIR / "Recall Bank/Bank_Imp topics and Questions.docx",
        "priority": 1, "type": "recall_bank",
        "topics": ["all"], "korky": True,
        "note": "Important topics and questions (docx)",
    },
    "recall-post-exam-notes": {
        "label": "Recall Bank — Post-Exam Notes",
        "path": KORKY_DIR / "Recall Bank/My post-exam notes.docx",
        "priority": 1, "type": "recall_bank",
        "topics": ["all"], "korky": True,
        "note": "Personal post-exam notes (docx)",
    },

    # ── BOOKS FOR FRCS PT.1 — previously unregistered ───────────────────
    "bank-neuro-surgery": {
        "label": "Bank — Neuro-Surgery MCQs (14pp)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Bank_7247248-Neuro-Surgery.pdf",
        "priority": 2, "type": "recall_bank",
        "topics": ["all"], "korky": True,
    },
    "bank-neurosurgery-mcqs": {
        "label": "Bank — Neurosurgery MCQs (5pp)",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Bank_173051952-Neurosurgery-MCQ-s.pdf",
        "priority": 2, "type": "recall_bank",
        "topics": ["all"], "korky": True,
    },
    "damirez-operative-anatomy": {
        "label": "Damirez/Fossett — Operative Neurosurgical Anatomy",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Damirez_Fossett_Operative_Neurosurgical_Anatomy_Thieme_Publishing.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["cranial-anatomy", "spinal-anatomy"], "korky": True,
    },
    "emergency-head-injury": {
        "label": "Emergency Surgery for Head Injury",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Emergency Surgery for Head Injury.pdf",
        "priority": 3, "type": "textbook",
        "topics": ["head-injury"], "korky": True,
    },
    "nader-cases": {
        "label": "Nader — Neurosurgery Cases Book",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Nader Neurosurgery cases book.docx",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "alleyne-self-assessment": {
        "label": "Alleyne/Woodall/Citow — Board Review Q&A Self-Assessment",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Neurosurgery-Board-Review-Questions-and-Answers-for-Self-Assessment Alleyne Jr-Woodall-Citow.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "oxford-neurological-surgery": {
        "label": "Oxford Textbook of Neurological Surgery",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Oxford Textbook of Neurological Surgery.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "schmidek-sweet": {
        "label": "Schmidek & Sweet — Operative Neurosurgical Techniques",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Schmidek & Sweet Operative Neurosurgical Techniques.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "shaya-practice-questions": {
        "label": "Shaya — Neurosurgery Practice Questions & Answers 2nd ed",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Shaya et al_Neurosurgery Practice Questions and Answers 2nd Edition.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },

    # ── ABERDEEN COURSE MATERIAL (priority 3) ─────────────────────────────
    "aberdeen-cerebral-physiology": {
        "label": "Aberdeen — Cerebral Physiology",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/cerebral physiology.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-icu", "neurophysiology"], "korky": True,
        "note": "CBF, ICP, autoregulation, CPP, Monroe-Kellie — core viva topic",
    },
    "aberdeen-cbf-graphs": {
        "label": "Aberdeen — CBF Graphs",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/cbf - graphs.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-icu", "neurophysiology"], "korky": True,
    },
    "aberdeen-tjones-revision": {
        "label": "TJones — Revision Notes (comprehensive)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/TJones_revision notes.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "Comprehensive trainee revision notes — mine section by section",
    },
    "aberdeen-tjones-exam": {
        "label": "TJones — Examination Notes",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Copy of TJones_examination notes.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "Clinical examination technique notes for viva",
    },
    "aberdeen-aeds": {
        "label": "Aberdeen — AEDs (antiepileptic drugs)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/AED_s.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["epilepsy-surgery"], "korky": True,
    },
    "aberdeen-ica-anatomy": {
        "label": "Aberdeen — ICA Anatomy",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/ICA anatomy.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["cranial-anatomy", "vascular-aneurysm"], "korky": True,
    },
    "aberdeen-tcd": {
        "label": "Aberdeen — TCD (transcranial Doppler)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/TCD.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology", "neuro-icu"], "korky": True,
    },
    "aberdeen-ncs": {
        "label": "Aberdeen — NCS (nerve conduction studies)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/NCS.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neurophysiology", "peripheral-nerve"], "korky": True,
    },
    "aberdeen-visual-fields": {
        "label": "Aberdeen — Visual Fields",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Visual fields.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology", "cranial-anatomy"], "korky": True,
    },
    "aberdeen-neuro-death": {
        "label": "Aberdeen — Neurological Determination of Death",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/neurological_determination_of_death_testing.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-icu", "ethics"], "korky": True,
    },
    "aberdeen-who-brain-tumour": {
        "label": "WHO Brain Tumour Classification 2016",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/WHO_Brain_Tumour_2016.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-onco-cranial", "neuropathology"], "korky": True,
        "note": "Molecular classification — complements NG99 content",
    },
    "aberdeen-carpal-tunnel": {
        "label": "Aberdeen — Carpal Tunnel Syndrome",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Carpal-Tunnel-Syndrome_1.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["peripheral-nerve"], "korky": True,
    },
    "aberdeen-rescueicp": {
        "label": "Aberdeen — Rescue ICP strategies",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/rescueicp.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-icu", "head-injury"], "korky": True,
    },
    "aberdeen-ulnar-c8t1": {
        "label": "Aberdeen — Ulnar vs C8/T1 Radiculopathy",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Ulnar v C8_T1 radiculopathy.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["peripheral-nerve", "degenerative-spine"], "korky": True,
    },
    "aberdeen-language": {
        "label": "Aberdeen — Language Organization (Mitch Berger)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/language organization_ Mitch Berger.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["cranial-anatomy", "neuro-onco-cranial"], "korky": True,
    },
    "aberdeen-eye-exam": {
        "label": "Aberdeen — Eye Examination",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/eye exam.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["cranial-anatomy", "neuroradiology"], "korky": True,
    },
    "aberdeen-la-toxicity": {
        "label": "Aberdeen — Local Anaesthetic Toxicity",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/la_toxicity.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuro-icu"], "korky": True,
    },

    # ── ABERDEEN — previously unregistered ─────────────────────────────────
    "aberdeen-exam-tips": {
        "label": "Aberdeen — Exam Tips and Tricks",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Exam tips and tricks.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
    },
    "aberdeen-tjones-exam-orig": {
        "label": "TJones — Examination Notes (original)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/TJones_examination notes.pdf",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "Original copy — compare with 'Copy of' version for differences",
    },
    "aberdeen-visual-fields-goldman": {
        "label": "Aberdeen — Visual Fields (Goldman)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Visual Fields (Goldman).pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology", "cranial-anatomy"], "korky": True,
    },
    "aberdeen-visual-fields-humphrey": {
        "label": "Aberdeen — Visual Fields (Humphrey)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Visual Fields (Humphrey).pdf",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology", "cranial-anatomy"], "korky": True,
    },
    "aberdeen-viva-q-jim": {
        "label": "Aberdeen — Viva Questions (Jim)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/Viva Q - Jim.docx",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
    },
    "aberdeen-brainstem": {
        "label": "Aberdeen — Brain Stem (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/brain stem.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["cranial-anatomy"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-frcs-sn-exam": {
        "label": "Aberdeen — FRCS(SN) EXAM (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/FRCS(SN) EXAM.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-frcs-clinicals": {
        "label": "Aberdeen — FRCS Clinicals (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/FRCS clinicals.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-frcs-close": {
        "label": "Aberdeen — FRCS Close (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/FRCS close.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-frcs-exam": {
        "label": "Aberdeen — FRCS Exam (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/FRCS exam.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["all"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-mri": {
        "label": "Aberdeen — MRI (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/MRI.pptx",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology"], "korky": True,
        "note": "PowerPoint — needs manual text extraction or conversion",
    },
    "aberdeen-radiology": {
        "label": "Aberdeen — Radiology FRCS (PowerPoint)",
        "path": KORKY_DIR / "Books for FRCS/Part 2 Prep - Aberdeen Course Material - JSW/RadiologyFRCS.ppt",
        "priority": 3, "type": "course_material",
        "topics": ["neuroradiology"], "korky": True,
        "note": "PowerPoint (.ppt) — needs manual text extraction or conversion",
    },

    # ── EXISTING GUIDELINES — NOT YET FULLY MINED ─────────────────────────
    "ng217-epilepsy": {
        "label": "NICE NG217 — Epilepsies (2022)",
        "path": APP_DIR / "guidelines/nice_ng217_epilepsy.pdf",
        "priority": 2, "type": "nice_guideline",
        "topics": ["epilepsy-surgery"], "korky": False,
        "note": "150pp — 0 cards mined from PDF so far. High priority.",
    },
    "rcp-stroke": {
        "label": "RCP/SIGN Stroke Guideline 2023",
        "path": APP_DIR / "guidelines/rcp_stroke_guideline_2023.pdf",
        "priority": 2, "type": "nice_guideline",
        "topics": ["vascular-ich", "carotid"], "korky": False,
        "note": "Large — only ICH/carotid partially covered. Thrombectomy, thrombolysis sections unmined.",
    },
    "ng127-neurological": {
        "label": "NICE NG127 — Suspected Neurological Conditions (2023)",
        "path": APP_DIR / "guidelines/nice_ng127_suspected_neurological.pdf",
        "priority": 3, "type": "nice_guideline",
        "topics": ["ethics"], "korky": False,
        "note": "86pp — basic content only",
    },

    # ── LARGE REFERENCE BOOKS (priority 4 — mine by chapter/section) ──────
    "greenberg": {
        "label": "Greenberg — Handbook of Neurosurgery 10th ed",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Greenberg's Handbook of Neurosurgery 10edition.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
        "note": "Primary reference. Mine chapter by chapter. Very large — approach by topic gap.",
    },
    "alleyne-board-review": {
        "label": "Alleyne/Citow — Neurosurgery Board Review 3E",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Alleyne Jr, Citow Thieme - Neurosurgery Board Review 3E.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
        "note": "Q&A format — mine topic sections that match weak areas",
    },
    "birinyi-board-prep": {
        "label": "Birinyi — Comprehensive Neurosurgery Board Prep",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Birinyi et al_The Comprehensive Neurosurgery Board Preparation Book.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "infographic-2025": {
        "label": "Infographic Guide to Neurosurgery 2025",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/The Infographic Guide to Neurosurgery 2025.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
        "note": "Visual summaries — mine for high-yield factual content per topic",
    },
    "harbaugh-knowledge-update": {
        "label": "Harbaugh — Neurosurgery Knowledge Update 1st ed",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/Harbaugh et al_Neurosurgery_Knowledge_Update_A_Comprehensive_Review_1st_Edition_medibos.blogspot.com.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["all"], "korky": True,
    },
    "spine-surgery-2019": {
        "label": "Spine Surgery 2019",
        "path": KORKY_DIR / "Books for FRCS/Books for FRCS Pt.1/2019_Book_SpineSurgery.pdf",
        "priority": 4, "type": "textbook",
        "topics": ["degenerative-spine", "spinal-trauma", "spinal-anatomy"], "korky": True,
    },
}


# ── CONTENT SOURCE ID MAP ────────────────────────────────────────────────────
# Maps mine.py source IDs → content.js SOURCES array IDs (for src_id field on cards)
CONTENT_SOURCE_MAP = {
    "paper-sivakumar-hyponatraemia": "k-sivakumar-1994",
    "paper-nascis3":                 "k-nascis3-1997",
    "paper-isuia":                   "k-isuia-1998",
    "paper-patchell":                "k-patchell-2005",
    "paper-stupp":                   "k-stupp-2005",
    "paper-sport":                   "k-sport-2006",
    "paper-berger":                  "k-berger-2008",
    "paper-isat":                    "k-isat-2009",
    "paper-santarius":               "k-santarius-2009",
    "paper-smith-abscess":           "k-smith-2009",
    "paper-hashimoto-desh":          "k-hashimoto-2010",
    "paper-kulkarni":                "k-kulkarni-2010",
    "paper-anand-mvd":               "k-anand-2011",
    "paper-decra":                   "k-decra-2011",
    "paper-kumpe-iih-stents":        "k-kumpe-2011",
    "paper-chesnut":                 "k-chesnut-2012",
    "paper-mautner-nf2":             "k-mautner-2010",
    "paper-couture-helmet":          "k-couture-2013",
    "paper-delwel-nph":              "k-delwel-2013",
    "paper-hu-epilepsy":             "k-hu-2013",
    "paper-jahangiri-redo-tss":      "k-jahangiri-2014",
    "paper-aruba":                   "k-aruba-2014",
    "paper-destiny2":                "k-destiny2-2014",
    "paper-liu-dbs":                 "k-liu-2014",
    # Guidelines
    "ng217-epilepsy":                "nice-ng217",
    "rcp-stroke":                    "stroke-guideline",
    "ng127-neurological":            "nice-ng127",
    "bank-neuro-surgery":            "k-bank-neuro-surgery",
    "bank-neurosurgery-mcqs":        "k-bank-mcqs",
    "aberdeen-neuro-death":          "aberdeen-neuro-death",
    "aberdeen-tcd":                  "aberdeen-tcd",
    "aberdeen-aeds":                 "aberdeen-aeds",
    "aberdeen-tjones-revision":      "aberdeen-tjones-revision",
    "aberdeen-tjones-exam":          "aberdeen-tjones-exam",
    # Future papers — add entry here when source is added to content.js SOURCES array
}

# ── MANIFEST ─────────────────────────────────────────────────────────────────
def load_manifest():
    if MANIFEST.exists():
        with open(MANIFEST) as f:
            return json.load(f)
    return {"sources": {}, "sessions": []}


def save_manifest(m):
    with open(MANIFEST, "w") as f:
        json.dump(m, f, indent=2)


def get_source_manifest(m, sid):
    if sid not in m["sources"]:
        m["sources"][sid] = {"pages_done": [], "cards_added": 0, "sessions": []}
    return m["sources"][sid]


# ── PDF UTILITIES ─────────────────────────────────────────────────────────────
def get_page_count(path):
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    except Exception:
        try:
            from pypdf import PdfReader
            return len(PdfReader(str(path)).pages)
        except Exception:
            return None


def extract_pdf_text(path, start_page=None, end_page=None):
    """Extract text from a PDF. Pages are 1-indexed."""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            total = len(pdf.pages)
            s = (start_page or 1) - 1
            e = min(end_page or total, total)
            chunks = []
            for i in range(s, e):
                page = pdf.pages[i]
                text = page.extract_text()
                if text and text.strip():
                    chunks.append(f"--- PAGE {i+1} ---\n{text.strip()}")
            return "\n\n".join(chunks), total
    except Exception as ex:
        return f"[ERROR extracting text: {ex}]", None


def extract_docx_text(path):
    try:
        from docx import Document
        doc = Document(str(path))
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(lines)
    except Exception as ex:
        return f"[ERROR extracting docx: {ex}]"


def extract_pptx_text(path):
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- SLIDE {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else "[No text content found in PowerPoint]"
    except ImportError:
        return "[python-pptx not installed: pip install python-pptx]"
    except Exception as ex:
        return f"[ERROR extracting pptx: {ex}]"


def extract_text(path, start_page=None, end_page=None):
    p = Path(path)
    if not p.exists():
        return f"[FILE NOT FOUND: {path}]", None
    if p.suffix.lower() == ".pdf":
        return extract_pdf_text(p, start_page, end_page)
    elif p.suffix.lower() in (".docx", ".doc"):
        return extract_docx_text(p), None
    elif p.suffix.lower() in (".pptx", ".ppt"):
        return extract_pptx_text(p), None
    elif p.suffix.lower() == ".txt":
        return p.read_text(errors="replace"), None
    else:
        return f"[Unsupported file type: {p.suffix}]", None


def parse_pages(page_str):
    """Parse '1-20' or '5' into (start, end) ints."""
    if not page_str:
        return None, None
    if "-" in page_str:
        parts = page_str.split("-")
        return int(parts[0]), int(parts[1])
    return int(page_str), int(page_str)


# ── CONTENT.JS PARSING ────────────────────────────────────────────────────────
def parse_content_js():
    """
    Returns dict: {topic_id: {cards: N, sba: N, korky_cards: N, korky_sba: N}}
    by scanning content.js with regex (avoids eval).
    """
    text = CONTENT.read_text()
    topics = {}

    # Find all topic blocks
    pattern = re.compile(r'"([a-z][a-z-]+)":\{src:"([^"]+)",c:\[')
    for m in pattern.finditer(text):
        tid = m.group(1)
        block_start = m.start()
        # Rough count of {q: (flashcards) and stem: (SBAs) within ~50k chars after
        block = text[block_start:block_start + 250000]
        # Find end of this topic block (next topic or end)
        next_topic = re.search(r'"[a-z][a-z-]+":\{src:', block[10:])
        if next_topic:
            block = block[:next_topic.start() + 10]

        cards  = len(re.findall(r'\{q:', block))
        sbas   = len(re.findall(r'\bstem:', block))
        kc     = len(re.findall(r'\{q:[^}]*korky:\s*true', block))
        ks     = len(re.findall(r'stem:[^}]*korky:\s*true', block))
        topics[tid] = {"cards": cards, "sba": sbas, "korky_cards": kc, "korky_sba": ks}

    return topics


def validate_content_js():
    """Check for common structural errors in content.js."""
    text = CONTENT.read_text()
    errors = []

    # Check every {q: has matching a:
    q_count  = len(re.findall(r'\{q:', text))
    a_count  = len(re.findall(r',a:"', text))
    if q_count != a_count:
        errors.append(f"Flashcard q/a mismatch: {q_count} q: vs {a_count} a:")

    # Check every stem: has opts: and ans:
    stem_count = len(re.findall(r'\bstem:', text))
    opts_count = len(re.findall(r'\bopts:', text))
    ans_count  = len(re.findall(r'\bans:', text))
    if stem_count != opts_count:
        errors.append(f"SBA stem/opts mismatch: {stem_count} stem: vs {opts_count} opts:")
    if stem_count != ans_count:
        errors.append(f"SBA stem/ans mismatch: {stem_count} stem: vs {ans_count} ans:")

    # Check ref fields exist on flashcards
    ref_count = len(re.findall(r',ref:"', text))
    if ref_count < q_count + stem_count:
        missing = (q_count + stem_count) - ref_count
        errors.append(f"{missing} cards/SBAs appear to be missing ref: fields")

    # Check refs have page numbers (contains 'p.' or a number after comma)
    refs = re.findall(r'ref:"([^"]+)"', text)
    no_page = [r for r in refs if not re.search(r'p\.?\s*\d|page|\d{4}|Table', r, re.I)]
    if no_page:
        errors.append(f"{len(no_page)} ref fields may be missing page/section numbers (sample: {no_page[:3]})")

    # Check LEARN object covers all TOPICS (only entries inside TOPICS=[...] block)
    topics_block = re.search(r'const TOPICS=\[(.*?)\];', text, re.S)
    if topics_block:
        topics_in_js = set(re.findall(r'\{id:"([a-z][a-z-]+)"', topics_block.group(1)))
        learn_in_js  = set(re.findall(r'"([a-z][a-z-]+)":\{src:', text))
        missing_learn = topics_in_js - learn_in_js
        if missing_learn:
            errors.append(f"LEARN entries missing for topics: {missing_learn}")

    return errors


# ── COMMANDS ──────────────────────────────────────────────────────────────────
def cmd_status(args):
    m = load_manifest()
    topic_stats = parse_content_js()

    print("\n" + "=" * 80)
    print("FRCS (SN) MINING STATUS")
    print("=" * 80)

    # Group by priority
    by_priority = {}
    for sid, s in SOURCES.items():
        p = s["priority"]
        by_priority.setdefault(p, []).append((sid, s))

    priority_labels = {
        1: "PRIORITY 1 — Recall Banks (real exam questions)",
        2: "PRIORITY 2 — Landmark Papers & Unmined Guidelines",
        3: "PRIORITY 3 — Aberdeen Course Material & Supporting PDFs",
        4: "PRIORITY 4 — Large Reference Textbooks (mine by chapter)",
    }

    total_sources = len(SOURCES)
    mined_sources = 0
    total_pages_est = 0
    pages_done = 0

    for p in sorted(by_priority.keys()):
        print(f"\n  {priority_labels[p]}")
        print("  " + "-" * 70)
        for sid, s in sorted(by_priority[p], key=lambda x: x[0]):
            sm = m["sources"].get(sid, {})
            done_pages = sm.get("pages_done", [])
            cards_added = sm.get("cards_added", 0)
            path = Path(s["path"])
            exists = path.exists()

            # Get page count
            pg = None
            if exists and path.suffix.lower() == ".pdf":
                pg = get_page_count(path)

            # Status icon
            if not exists:
                status = "MISSING"
                icon = "?"
            elif not done_pages:
                status = "not started"
                icon = " "
            elif pg and sum(e - st + 1 for st, e in done_pages) >= pg * 0.95:
                status = "COMPLETE"
                icon = "+"
                mined_sources += 1
            else:
                status = "in progress"
                icon = "~"

            pg_str = f"{pg}pp" if pg else ("docx" if path.suffix == ".docx" else "?pp")
            cards_str = f"{cards_added}c" if cards_added else "0c"
            korky_tag = "[K]" if s["korky"] else "   "

            print(f"  [{icon}] {korky_tag} {sid:<35} {pg_str:<7} {cards_str:<6}  {s['label'][:45]}")
            if s.get("note"):
                print(f"           {'':37}  NOTE: {s['note'][:60]}")

            if pg:
                total_pages_est += pg
                pages_done += sum(e - st + 1 for st, e in done_pages) if done_pages else 0

    # Summary
    topic_total = sum(v["cards"] for v in topic_stats.values())
    topic_korky = sum(v["korky_cards"] for v in topic_stats.values())
    sba_total   = sum(v["sba"] for v in topic_stats.values())
    sba_korky   = sum(v["korky_sba"] for v in topic_stats.values())

    print("\n" + "=" * 80)
    print(f"  Sources: {mined_sources}/{total_sources} fully mined")
    if total_pages_est:
        pct = round(pages_done / total_pages_est * 100) if total_pages_est else 0
        print(f"  Pages:   {pages_done}/{total_pages_est} ({pct}%)")
    print(f"  Cards:   {topic_total} flashcards ({topic_korky} Korky) | {sba_total} SBAs ({sba_korky} Korky)")
    print(f"\n  [K] = Korky folder source   [+] = complete   [~] = in progress   [ ] = not started")
    print("=" * 80 + "\n")


    # Mining order — by folder, systematically
MINING_ORDER = [
    # ── Folder 1: Recall Bank (docx files — PDFs already parsed into Recall tab) ──
    "recall-imp-topics", "recall-post-exam-notes",
    # ── Folder 2: Key Papers for FRCS (25 papers, alphabetical by year) ──
    "paper-sivakumar-hyponatraemia", "paper-casey-myelopathy", "paper-nascis3",
    "paper-isuia", "paper-patchell", "paper-stupp", "paper-sport",
    "paper-berger", "paper-isat", "paper-santarius", "paper-smith-abscess",
    "paper-hashimoto-desh", "paper-kulkarni", "paper-anand-mvd", "paper-decra",
    "paper-kumpe-iih-stents", "paper-chesnut", "paper-mautner-nf2",
    "paper-couture-helmet", "paper-delwel-nph", "paper-hu-epilepsy",
    "paper-jahangiri-redo-tss", "paper-destiny2", "paper-liu-dbs", "paper-aruba",
    # ── Folder 3: Books for FRCS Pt.1 (top-level files, excl. Key Papers subfolder) ──
    "bank-neuro-surgery", "bank-neurosurgery-mcqs",
    "emergency-head-injury", "infographic-2025",
    "nader-cases", "alleyne-board-review", "alleyne-self-assessment",
    "birinyi-board-prep", "shaya-practice-questions",
    "damirez-operative-anatomy", "spine-surgery-2019",
    "harbaugh-knowledge-update", "oxford-neurological-surgery",
    "schmidek-sweet", "greenberg",
    # ── Folder 4: Part 2 Prep - Aberdeen Course Material ──
    "aberdeen-cerebral-physiology", "aberdeen-cbf-graphs", "aberdeen-rescueicp",
    "aberdeen-aeds", "aberdeen-ica-anatomy", "aberdeen-tcd",
    "aberdeen-neuro-death", "aberdeen-la-toxicity",
    "aberdeen-language", "aberdeen-visual-fields",
    "aberdeen-visual-fields-goldman", "aberdeen-visual-fields-humphrey",
    "aberdeen-who-brain-tumour", "aberdeen-ulnar-c8t1",
    "aberdeen-carpal-tunnel", "aberdeen-ncs",
    "aberdeen-eye-exam", "aberdeen-exam-tips",
    "aberdeen-tjones-revision", "aberdeen-tjones-exam", "aberdeen-tjones-exam-orig",
    "aberdeen-viva-q-jim",
    "aberdeen-brainstem", "aberdeen-frcs-sn-exam",
    "aberdeen-frcs-clinicals", "aberdeen-frcs-close",
    "aberdeen-frcs-exam", "aberdeen-mri", "aberdeen-radiology",
]


def _source_folder(sid):
    if sid.startswith("recall"):
        return "Recall Bank"
    if sid.startswith("paper-"):
        return "Key Papers"
    if sid.startswith("aberdeen"):
        return "Aberdeen"
    return "Books Pt.1"


def cmd_next(args):
    m = load_manifest()

    # Count overall progress
    total_sources = len(MINING_ORDER)
    done_sources = 0
    total_pages = 0
    done_pages_count = 0

    for sid in MINING_ORDER:
        s = SOURCES[sid]
        path = Path(s["path"])
        pg = get_page_count(path) if path.exists() and path.suffix == ".pdf" else None
        if pg:
            total_pages += pg
        sm = m["sources"].get(sid, {})
        dp = sm.get("pages_done", [])
        mined = sum(e - st + 1 for st, e in dp) if dp else 0
        if pg:
            done_pages_count += min(mined, pg)
        if pg and mined >= pg * 0.95:
            done_sources += 1

    pct = round(done_pages_count / total_pages * 100) if total_pages else 0
    print(f"\n  Mining progress: {done_sources}/{total_sources} sources, "
          f"{done_pages_count}/{total_pages} pages ({pct}%)\n")

    # Find next unmined or partially mined source in order
    shown = 0
    for sid in MINING_ORDER:
        if shown >= 3:
            break
        s = SOURCES[sid]
        path = Path(s["path"])
        if not path.exists():
            continue
        pg = get_page_count(path) if path.suffix == ".pdf" else None

        sm = m["sources"].get(sid, {})
        done_pages = sm.get("pages_done", [])

        if not done_pages:
            chunk_end = min(20, pg) if pg else 20
            folder = _source_folder(sid)
            marker = ">>>" if shown == 0 else "   "
            print(f"  {marker} [{folder}]  {sid}  ({pg}pp)")
            print(f"      {s['label']}")
            print(f"      python3 mine.py extract {sid} 1-{chunk_end}")
            if s.get("note"):
                print(f"      Note: {s['note']}")
            print()
            shown += 1
        elif pg:
            mined_flat = set()
            for st, e in done_pages:
                mined_flat.update(range(st, e + 1))
            unmined = [i for i in range(1, pg + 1) if i not in mined_flat]
            if unmined:
                start = unmined[0]
                end = min(start + 19, pg)
                folder = _source_folder(sid)
                marker = ">>>" if shown == 0 else "   "
                print(f"  {marker} [{folder}]  {sid}  (resuming — {len(mined_flat)}/{pg} pages done)")
                print(f"      {s['label']}")
                print(f"      python3 mine.py extract {sid} {start}-{end}")
                print()
                shown += 1

    if shown == 0:
        print("  All sources fully mined! Nothing left to do.")


def cmd_extract(args):
    if len(args) < 1:
        print("Usage: python3 mine.py extract <source_id> [pages, e.g. 1-20]")
        sys.exit(1)

    sid = args[0]
    if sid not in SOURCES:
        # Try partial match
        matches = [k for k in SOURCES if k.startswith(sid)]
        if len(matches) == 1:
            sid = matches[0]
        elif len(matches) > 1:
            print(f"Ambiguous ID '{sid}'. Matches: {matches}")
            sys.exit(1)
        else:
            print(f"Unknown source: '{sid}'. Run 'python3 mine.py status' to see all IDs.")
            sys.exit(1)

    s = SOURCES[sid]
    page_str = args[1] if len(args) > 1 else None
    start, end = parse_pages(page_str)

    path = Path(s["path"])
    if not path.exists():
        print(f"[FILE NOT FOUND] {path}")
        sys.exit(1)

    print(f"\n{'=' * 80}")
    print(f"SOURCE:  {s['label']}")
    print(f"ID:      {sid}")
    print(f"TYPE:    {s['type']}  |  KORKY: {s['korky']}")
    print(f"TOPICS:  {', '.join(s['topics'])}")
    if s.get("note"):
        print(f"NOTE:    {s['note']}")
    print(f"{'=' * 80}\n")

    text, total_pages = extract_text(path, start, end)

    if total_pages:
        print(f"Pages {start or 1}–{end or total_pages} of {total_pages}\n")

    print(text)

    print(f"\n{'=' * 80}")
    print(f"MINING INSTRUCTIONS FOR CLAUDE:")
    print(f"  Source ID : {sid}")
    print(f"  Pages     : {start or 1}–{end or total_pages or '?'}")
    print(f"  korky tag : {str(s['korky']).lower()}")
    print(f"  Topics    : {', '.join(s['topics'])}")
    print(f"  Ref format: e.g. '{s['label'].split(' — ')[0]}, p.{start or 1}'")
    content_id = CONTENT_SOURCE_MAP.get(sid)
    if content_id:
        print(f"  src_id tag: {content_id}  (add src_id:\"{content_id}\" to every card)")
    else:
        print(f"  src_id tag: NOT MAPPED — add this source to CONTENT_SOURCE_MAP in mine.py first")
    print(f"\n  After adding cards to content.js, mark as done:")
    print(f"  python3 mine.py done {sid} {start or 1}-{end or total_pages or '?'} <N_cards_added>")
    print(f"{'=' * 80}\n")


def cmd_done(args):
    if len(args) < 3:
        print("Usage: python3 mine.py done <source_id> <pages e.g. 1-20> <N_cards_added>")
        sys.exit(1)

    sid, page_str, n_str = args[0], args[1], args[2]
    if sid not in SOURCES:
        print(f"Unknown source: {sid}")
        sys.exit(1)

    start, end = parse_pages(page_str)
    n_cards = int(n_str)

    m = load_manifest()
    sm = get_source_manifest(m, sid)

    # Merge page range
    sm["pages_done"].append([start, end])
    sm["cards_added"] = sm.get("cards_added", 0) + n_cards
    sm["sessions"].append({
        "date": str(date.today()),
        "pages": f"{start}-{end}",
        "cards": n_cards,
    })

    save_manifest(m)
    print(f"\n  Marked {sid} pages {start}-{end} as done. +{n_cards} cards.")
    print(f"  Total cards from this source: {sm['cards_added']}\n")


def cmd_validate(args):
    print("\n  Validating content.js...\n")
    errors = validate_content_js()
    if not errors:
        print("  No structural errors found.\n")
    else:
        for e in errors:
            print(f"  ERROR: {e}")
        print()


def cmd_stats(args):
    topic_stats = parse_content_js()
    m = load_manifest()

    print("\n" + "=" * 80)
    print("CARD STATISTICS BY TOPIC")
    print("=" * 80)
    print(f"  {'Topic':<35} {'Cards':>6} {'SBAs':>5} {'Korky-C':>8} {'Korky-Q':>8}")
    print("  " + "-" * 66)

    total_c = total_s = total_kc = total_ks = 0
    for tid, v in sorted(topic_stats.items()):
        print(f"  {tid:<35} {v['cards']:>6} {v['sba']:>5} {v['korky_cards']:>8} {v['korky_sba']:>8}")
        total_c  += v["cards"]
        total_s  += v["sba"]
        total_kc += v["korky_cards"]
        total_ks += v["korky_sba"]

    print("  " + "-" * 66)
    print(f"  {'TOTAL':<35} {total_c:>6} {total_s:>5} {total_kc:>8} {total_ks:>8}")

    print("\n  Mining sessions:")
    for sess in m.get("sessions", []):
        print(f"    {sess.get('date','?')}  {sess.get('source','?'):<30} pages {sess.get('pages','?')}  +{sess.get('cards',0)} cards")
    if not m.get("sessions"):
        print("    No sessions recorded yet.")
    print()


def cmd_refs(args):
    """List all ref strings — useful for spotting cards missing page numbers."""
    text = CONTENT.read_text()
    refs = re.findall(r'ref:"([^"]+)"', text)
    no_page = [r for r in refs if not re.search(r'p\.?\s*\d|page|\d{4}|Table', r, re.I)]
    print(f"\n  Total ref fields: {len(refs)}")
    print(f"  Refs possibly missing page/section: {len(no_page)}")
    if no_page:
        print("\n  These refs have no page number:")
        for r in sorted(set(no_page)):
            print(f"    {r}")
    print()


# ── MAIN ──────────────────────────────────────────────────────────────────────
COMMANDS = {
    "status":   cmd_status,
    "next":     cmd_next,
    "extract":  cmd_extract,
    "done":     cmd_done,
    "validate": cmd_validate,
    "stats":    cmd_stats,
    "refs":     cmd_refs,
}

if __name__ == "__main__":
    if len(sys.argv) < 2 or sys.argv[1] not in COMMANDS:
        print(__doc__)
        print("  Available commands:", ", ".join(COMMANDS))
        sys.exit(0)
    COMMANDS[sys.argv[1]](sys.argv[2:])
